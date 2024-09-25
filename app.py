from flask import Flask, render_template, request, redirect, url_for, send_file
import mysql.connector
import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import io
import hashlib
from datetime import datetime
import plotly.express as px
import plotly.io as pio

app = Flask(__name__)

# Configuração do banco de dados
db_config = {
    'user': 'root',
    'password': '',
    'host': 'localhost',
    'database': 'tca'
}

db = mysql.connector.connect(**db_config)

def generate_unique_short_code():
    """Gera um código curto único."""
    while True:
        short_code = hashlib.md5(datetime.now().strftime('%Y-%m-%d %H:%M:%S').encode()).hexdigest()[:6]
        cursor = db.cursor()
        cursor.execute("SELECT COUNT(*) FROM urls WHERE short_code = %s", (short_code,))
        if cursor.fetchone()[0] == 0:
            cursor.close()
            return short_code
        cursor.close()

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        original_url = request.form["original_url"]
        short_code = generate_unique_short_code()

        cursor = db.cursor()
        cursor.execute("INSERT INTO urls (original_url, short_code, created_at, click_count) VALUES (%s, %s, %s, %s)",
                       (original_url, short_code, datetime.now(), 0))
        db.commit()
        cursor.close()

        short_url = url_for("redirect_url", short_code=short_code, _external=True)
        return render_template("index.html", short_url=short_url)

    return render_template("index.html")

@app.route("/<short_code>")
def redirect_url(short_code):
    cursor = db.cursor(dictionary=True)
    cursor.execute("SELECT original_url FROM urls WHERE short_code = %s", (short_code,))
    url_data = cursor.fetchone()

    if url_data:
        client_ip = request.remote_addr
        cursor.execute("""
            UPDATE urls 
            SET last_click_at = %s, last_click_ip = %s, click_count = click_count + 1 
            WHERE short_code = %s
        """, (datetime.now(), client_ip, short_code))
        db.commit()
        cursor.close()
        return redirect(url_data['original_url'])
    else:
        cursor.close()
        return "URL não encontrada", 404

@app.route("/urls")
def show_urls():
    cursor = db.cursor(dictionary=True)
    cursor.execute("SELECT * FROM urls")
    urls = cursor.fetchall()
    cursor.close()
    return render_template("urls.html", urls=urls)

@app.route("/charts")
def charts():
    cursor = db.cursor(dictionary=True)
    cursor.execute("SELECT short_code, original_url, click_count, created_at, last_click_at, last_click_ip FROM urls")
    urls = cursor.fetchall()
    cursor.close()

    # Preparar os dados para o gráfico
    df = pd.DataFrame(urls)
    fig = px.bar(df, x='short_code', y='click_count', labels={'short_code': 'Short Code', 'click_count': 'Clicks'}, title='URL Clicks')

    # Converter gráfico para HTML
    chart_html = fig.to_html(full_html=False)

    return render_template("charts.html", url_data=urls, chart_html=chart_html)

@app.route("/download/<file_type>")
def download_report(file_type):
    cursor = db.cursor(dictionary=True)
    cursor.execute("SELECT short_code, original_url, click_count, created_at, last_click_at, last_click_ip FROM urls")
    urls = cursor.fetchall()
    cursor.close()

    if file_type == 'pptx':
        prs = Presentation()
        
        # Gerar o gráfico e adicionar ao PowerPoint
        df = pd.DataFrame(urls)
        fig = px.bar(df, x='short_code', y='click_count', labels={'short_code': 'Short Code', 'click_count': 'Clicks'}, title='URL Clicks')
        img_bytes = pio.to_image(fig, format='png')
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "URL Clicks"

        # Adicionar gráfico ao slide
        image_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), width=Inches(8))

        # Adicionar informações das URLs
        for url in urls:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"Short Code: {url['short_code']}"

            content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            text_frame = content.text_frame
            p = text_frame.add_paragraph()
            p.text = (
                f"Original URL: {url['original_url']}\n"
                f"Clicks: {url['click_count']}\n"
                f"Created At: {url['created_at'].strftime('%Y-%m-%d %H:%M:%S')}\n"
                f"Last Click At: {url['last_click_at'].strftime('%Y-%m-%d %H:%M:%S') if url['last_click_at'] else 'Never'}\n"
                f"Last Click IP: {url['last_click_ip'] if url['last_click_ip'] else 'N/A'}"
            )

        # Salvar o arquivo pptx em memória
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return send_file(file_stream, as_attachment=True, download_name="report.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    elif file_type == 'docx':
        doc = Document()
        doc.add_heading('URL Report', 0)

        for url in urls:
            doc.add_paragraph(f"Short Code: {url['short_code']}")
            doc.add_paragraph(f"Original URL: {url['original_url']}")
            doc.add_paragraph(f"Clicks: {url['click_count']}")
            doc.add_paragraph(f"Created At: {url['created_at'].strftime('%Y-%m-%d %H:%M:%S')}")
            doc.add_paragraph(f"Last Click At: {url['last_click_at'].strftime('%Y-%m-%d %H:%M:%S') if url['last_click_at'] else 'Never'}")
            doc.add_paragraph(f"Last Click IP: {url['last_click_ip'] if url['last_click_ip'] else 'N/A'}")
            doc.add_paragraph('')  # Linha em branco para separar as URLs

        # Salvar o arquivo docx em memória
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return send_file(file_stream, as_attachment=True, download_name="report.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    elif file_type == 'xlsx':
        df = pd.DataFrame(urls)
        df['created_at'] = df['created_at'].apply(lambda x: x.strftime('%Y-%m-%d %H:%M:%S'))
        df['last_click_at'] = df['last_click_at'].apply(lambda x: x.strftime('%Y-%m-%d %H:%M:%S') if x else 'Never')
        df['last_click_ip'] = df['last_click_ip'].fillna('N/A')

        file_stream = io.BytesIO()
        df.to_excel(file_stream, index=False)
        file_stream.seek(0)

        return send_file(file_stream, as_attachment=True, download_name="report.xlsx", mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    else:
        return "Tipo de arquivo não suportado", 400

if __name__ == "__main__":
    app.run(debug=True)
